"""
PPTX/PPT format handler — slide extraction and reconstruction via python-pptx.

Ingest:
  Each slide → H2 section. Extracts titles, body text, tables, images, speaker notes.
  Charts rendered as images (libreoffice mode) or placeholders. SmartArt detected
  with text extraction and a warning. Grouped shapes recursed where possible.
  .ppt files are first converted to .pptx via LibreOffice headless.

Export:
  Rebuilds PPTX from H2-delimited sections.
  Tier 2: apply slide dimensions, font properties from sidecar.
  Tier 3: if original .pptx exists and hash match ≥ 80%, patch text in-place.
"""

import hashlib
import io
import time
from pathlib import Path
from typing import Any

import structlog

from formats.base import FormatHandler, register_handler
from core.document_model import (
    DocumentModel,
    DocumentMetadata,
    Element,
    ElementType,
    ImageData,
    compute_content_hash,
)

log = structlog.get_logger(__name__)


@register_handler
class PptxHandler(FormatHandler):
    EXTENSIONS = ["pptx", "ppt", "pptm"]

    # ── Ingest ────────────────────────────────────────────────────────────────

    def ingest(self, file_path: Path) -> DocumentModel:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        file_path = Path(file_path)
        t_start = time.perf_counter()
        log.info("handler_ingest_start", filename=file_path.name, format="pptx")
        _tmp_pptx: Path | None = None

        if file_path.suffix.lower() == ".ppt":
            from core.libreoffice_helper import convert_with_libreoffice

            file_path = convert_with_libreoffice(file_path, "pptx")
            _tmp_pptx = file_path

        try:
            # Read chart extraction preference synchronously (same pattern
            # as database_handler._get_sample_rows_limit)
            self._chart_mode = self._read_chart_mode_pref()
            self._slide_images: dict[int, Any] = {}  # lazy-loaded per ingest

            model = DocumentModel()
            model.metadata = DocumentMetadata(
                source_file=file_path.name,
                source_format="pptx",
            )

            prs = Presentation(str(file_path))
            slide_count = len(prs.slides)
            model.metadata.page_count = slide_count

            # Slide dimensions in EMU (needed for chart crop coordinate conversion)
            self._slide_width_emu = prs.slide_width
            self._slide_height_emu = prs.slide_height

            # Pre-render slides if libreoffice chart mode is active and file
            # contains at least one chart shape
            if self._chart_mode == "libreoffice":
                has_charts = any(
                    hasattr(shape, "has_chart") and shape.has_chart
                    for slide in prs.slides
                    for shape in slide.shapes
                )
                if has_charts:
                    self._slide_images = self._render_slides_via_libreoffice(file_path)

            for idx, slide in enumerate(prs.slides):
                # Add horizontal rule between slides (not before first)
                if idx > 0:
                    model.add_element(Element(type=ElementType.HORIZONTAL_RULE, content=""))

                # Slide title
                title_text = self._get_slide_title(slide, idx + 1)
                model.add_element(
                    Element(type=ElementType.HEADING, content=title_text, level=2)
                )

                # Walk shapes
                for shape in slide.shapes:
                    self._process_shape(shape, model, file_path, idx + 1)

                # Speaker notes
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        model.add_element(
                            Element(
                                type=ElementType.BLOCKQUOTE,
                                content=f"Speaker Notes: {notes_text}",
                            )
                        )

            duration_ms = int((time.perf_counter() - t_start) * 1000)
            log.info(
                "handler_ingest_complete",
                filename=file_path.name,
                element_count=len(model.elements),
                duration_ms=duration_ms,
            )
            return model
        finally:
            self._slide_images = {}  # free memory
            if _tmp_pptx and _tmp_pptx.exists():
                _tmp_pptx.unlink(missing_ok=True)

    # ── Preference helpers ───────────────────────────────────────────────────

    @staticmethod
    def _read_chart_mode_pref() -> str:
        """Return the pptx_chart_extraction_mode preference.

        Prefers the in-memory preferences cache (no DB I/O). On a cold
        cache, falls back to a single synchronous sqlite read and warms
        the cache so subsequent conversions hit memory only. This avoids
        the "open a raw sqlite3 connection on every PPTX ingest" pattern
        that bypassed the connection pool.
        """
        try:
            from core.preferences_cache import peek_cached_preference, _cache
            import time as _time

            value = peek_cached_preference("pptx_chart_extraction_mode")
            if value in ("placeholder", "libreoffice"):
                return value

            # Cold cache: one-time sync read, then populate the cache so
            # all later PPTX conversions are cache-only.
            import sqlite3
            from core.database import get_db_path

            conn = sqlite3.connect(get_db_path())
            try:
                row = conn.execute(
                    "SELECT value FROM user_preferences WHERE key = ?",
                    ("pptx_chart_extraction_mode",),
                ).fetchone()
            finally:
                conn.close()

            resolved = row[0] if row and row[0] in ("placeholder", "libreoffice") else "placeholder"
            _cache["pptx_chart_extraction_mode"] = (resolved, _time.time() + 300)
            return resolved
        except Exception:
            return "placeholder"

    # ── LibreOffice slide rendering ──────────────────────────────────────────

    def _render_slides_via_libreoffice(self, file_path: Path) -> dict[int, Any]:
        """Convert PPTX to PDF via LibreOffice, render each page to a PIL Image.

        Returns dict mapping slide_num (1-based) to PIL Image.
        """
        try:
            from core.libreoffice_helper import convert_with_libreoffice
            import fitz  # PyMuPDF

            pdf_path = convert_with_libreoffice(file_path, "pdf", timeout=60)
            try:
                doc = fitz.open(str(pdf_path))
                images: dict[int, Any] = {}
                for i in range(len(doc)):
                    page = doc[i]
                    # Render at 2x for decent quality
                    pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                    from PIL import Image

                    img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
                    images[i + 1] = img  # 1-based slide numbers
                doc.close()
                return images
            finally:
                pdf_path.unlink(missing_ok=True)
        except Exception as exc:
            log.warning("pptx.libreoffice_render_failed", error=str(exc))
            return {}

    def _crop_shape_region(
        self, slide_img: Any, shape: Any, slide_width_emu: int, slide_height_emu: int
    ) -> Any:
        """Crop a region from a rendered slide image matching the shape's bounds.

        Returns a PIL Image or None on failure.
        """
        try:
            img_w, img_h = slide_img.size
            # EMU to pixel conversion
            scale_x = img_w / slide_width_emu
            scale_y = img_h / slide_height_emu

            left = int(shape.left * scale_x)
            top = int(shape.top * scale_y)
            right = int((shape.left + shape.width) * scale_x)
            bottom = int((shape.top + shape.height) * scale_y)

            # Clamp to image bounds
            left = max(0, left)
            top = max(0, top)
            right = min(img_w, right)
            bottom = min(img_h, bottom)

            if right <= left or bottom <= top:
                return None

            return slide_img.crop((left, top, right, bottom))
        except Exception as exc:
            log.debug("pptx.chart_crop_failed", error=str(exc))
            return None

    def _get_slide_title(self, slide: Any, slide_num: int) -> str:
        """Extract slide title from the title placeholder, or generate one."""
        if slide.shapes.title and slide.shapes.title.text.strip():
            return slide.shapes.title.text.strip()

        # Try to find any placeholder that looks like a title
        for shape in slide.placeholders:
            if hasattr(shape, "text") and shape.text.strip():
                if shape.placeholder_format.idx in (0, 1):  # Title or subtitle
                    return shape.text.strip()

        return f"Slide {slide_num}"

    def _process_shape(
        self, shape: Any, model: DocumentModel, file_path: Path, slide_num: int
    ) -> None:
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        # Skip title placeholder (already extracted)
        try:
            if shape.placeholder_format is not None:
                if shape.placeholder_format.idx == 0:  # Title
                    return
        except (ValueError, AttributeError):
            pass  # Not a placeholder — continue processing

        # Table
        if shape.has_table:
            rows: list[list[str]] = []
            for row in shape.table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(cells)
            if rows:
                model.add_element(Element(type=ElementType.TABLE, content=rows))
            return

        # Image / Picture
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            self._extract_shape_image(shape, model, slide_num)
            return

        # Text frame
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = self._extract_paragraph_text(para)
                if text.strip():
                    model.add_element(
                        Element(type=ElementType.PARAGRAPH, content=text)
                    )
            return

        # Group shape — recurse (with SmartArt detection)
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            # Check if this group is actually SmartArt
            shape_xml = shape._element.xml if hasattr(shape._element, "xml") else ""
            is_smartart = "dgm:relIds" in shape_xml or "smartArt" in shape_xml.lower()
            if is_smartart:
                model.warnings.append(
                    f"Slide {slide_num}: SmartArt not fully extractable"
                )
            # Recurse into children to get text content regardless
            try:
                for child_shape in shape.shapes:
                    self._process_shape(child_shape, model, file_path, slide_num)
            except Exception as exc:
                log.debug("pptx.group_recurse_failed", error=str(exc))
            return

        # Chart — render via LibreOffice or fall back to placeholder
        if hasattr(shape, "has_chart") and shape.has_chart:
            chart_title = ""
            try:
                chart_title = shape.chart.chart_title.text_frame.text
            except Exception:
                pass

            # Attempt LibreOffice rendering when enabled
            if self._chart_mode == "libreoffice" and self._slide_images:
                slide_img = self._slide_images.get(slide_num)
                if slide_img:
                    chart_img = self._crop_shape_region(
                        slide_img, shape,
                        self._slide_width_emu, self._slide_height_emu,
                    )
                    if chart_img is not None:
                        buf = io.BytesIO()
                        chart_img.save(buf, format="PNG")
                        img_bytes = buf.getvalue()
                        img_hash = hashlib.md5(img_bytes).hexdigest()[:12]
                        img_name = f"chart_{slide_num}_{img_hash}.png"
                        alt = f"Chart: {chart_title}" if chart_title else "Chart"
                        model.images[img_name] = ImageData(
                            data=img_bytes,
                            original_format="png",
                            width=chart_img.width,
                            height=chart_img.height,
                            alt_text=alt,
                        )
                        model.add_element(Element(
                            type=ElementType.IMAGE,
                            content=img_name,
                            attributes={"chart": True, "rendered_by": "libreoffice"},
                        ))
                        return

            # Fallback: placeholder mode
            model.add_element(
                Element(
                    type=ElementType.PARAGRAPH,
                    content=f"[Chart: {chart_title}]" if chart_title else "[Chart]",
                    attributes={"chart": True},
                )
            )
            model.warnings.append(f"Slide {slide_num}: chart not fully extractable")
            return

        # Embedded media
        if shape.shape_type in (
            MSO_SHAPE_TYPE.MEDIA,
            getattr(MSO_SHAPE_TYPE, "LINKED_OLE_OBJECT", -99),
            getattr(MSO_SHAPE_TYPE, "EMBEDDED_OLE_OBJECT", -99),
        ):
            name = getattr(shape, "name", "media")
            model.add_element(
                Element(
                    type=ElementType.RAW_HTML,
                    content=f"<!-- Embedded media: {name} -->",
                )
            )
            return

    def _extract_paragraph_text(self, para: Any) -> str:
        """Extract paragraph text with inline formatting markers."""
        parts: list[str] = []
        for run in para.runs:
            text = run.text
            if not text:
                continue
            if run.font.bold and run.font.italic:
                parts.append(f"***{text}***")
            elif run.font.bold:
                parts.append(f"**{text}**")
            elif run.font.italic:
                parts.append(f"*{text}*")
            else:
                parts.append(text)
        return "".join(parts)

    def _extract_shape_image(self, shape: Any, model: DocumentModel, slide_num: int) -> None:
        from core.image_handler import extract_image

        try:
            blob = shape.image.blob
            content_type = shape.image.content_type or "image/png"
            fmt = content_type.split("/")[-1]
            hash_name, png_data, meta = extract_image(blob, fmt)
            model.images[hash_name] = ImageData(
                data=png_data,
                original_format=fmt,
                width=meta.get("width"),
                height=meta.get("height"),
            )
            model.add_element(
                Element(
                    type=ElementType.IMAGE,
                    content=f"assets/{hash_name}",
                    attributes={"slide": slide_num},
                )
            )
        except Exception as exc:
            log.debug("pptx.image_extract_failed", slide=slide_num, error=str(exc))

    # ── Export ─────────────────────────────────────────────────────────────────

    def export(
        self,
        model: DocumentModel,
        output_path: Path,
        sidecar: dict[str, Any] | None = None,
        original_path: Path | None = None,
    ) -> None:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu

        tier = 3 if (original_path and original_path.exists() and sidecar) else (2 if sidecar else 1)
        t_start = time.perf_counter()
        log.info("handler_export_start", filename=output_path.name, target_format="pptx", tier=tier)

        # Tier 3: patch original
        if original_path and original_path.exists() and sidecar:
            if self._try_tier3_export(model, output_path, original_path, sidecar):
                duration_ms = int((time.perf_counter() - t_start) * 1000)
                log.info("handler_export_complete", filename=output_path.name, output_path=str(output_path), duration_ms=duration_ms)
                return

        prs = Presentation()

        # Apply slide dimensions from sidecar
        if sidecar:
            doc_level = sidecar.get("document_level", {})
            sw = doc_level.get("slide_width")
            sh = doc_level.get("slide_height")
            if sw and sh:
                prs.slide_width = int(sw)
                prs.slide_height = int(sh)

        # Parse model into slide sections (split on H2 boundaries)
        sections = self._split_into_slide_sections(model)

        for section in sections:
            self._add_slide(prs, section, model, sidecar)

        prs.save(str(output_path))
        duration_ms = int((time.perf_counter() - t_start) * 1000)
        log.info("handler_export_complete", filename=output_path.name, output_path=str(output_path), duration_ms=duration_ms)

    def _split_into_slide_sections(
        self, model: DocumentModel
    ) -> list[list[Element]]:
        """Split elements into per-slide groups at H2 boundaries."""
        sections: list[list[Element]] = []
        current: list[Element] = []

        for elem in model.elements:
            if elem.type == ElementType.HORIZONTAL_RULE and current:
                sections.append(current)
                current = []
                continue
            if elem.type == ElementType.HEADING and elem.level == 2 and current:
                sections.append(current)
                current = []
            current.append(elem)

        if current:
            sections.append(current)

        return sections

    def _add_slide(
        self,
        prs: Any,
        elements: list[Element],
        model: DocumentModel,
        sidecar: dict[str, Any] | None,
    ) -> None:
        from pptx.util import Inches, Pt, Emu

        # Determine layout
        layout_idx = 1  # Title and Content
        if sidecar and elements:
            h = compute_content_hash(elements[0].content)
            elem_style = sidecar.get("elements", {}).get(h, {})
            li = elem_style.get("layout_index")
            if li is not None:
                layout_idx = li

        try:
            layout = prs.slide_layouts[layout_idx]
        except (IndexError, KeyError):
            layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]

        slide = prs.slides.add_slide(layout)

        title_set = False
        body_parts: list[str] = []
        notes_text = ""

        for elem in elements:
            if elem.type == ElementType.HEADING and elem.level == 2 and not title_set:
                if slide.shapes.title:
                    slide.shapes.title.text = str(elem.content)
                title_set = True
                continue

            if elem.type == ElementType.BLOCKQUOTE:
                text = str(elem.content)
                if text.startswith("Speaker Notes:"):
                    notes_text = text[len("Speaker Notes:"):].strip()
                    continue
                body_parts.append(text)
                continue

            if elem.type == ElementType.TABLE:
                self._add_table_to_slide(slide, elem.content)
                continue

            if elem.type == ElementType.IMAGE:
                self._add_image_to_slide(slide, elem, model)
                continue

            if elem.type in (ElementType.PAGE_BREAK, ElementType.HORIZONTAL_RULE):
                continue

            if elem.type == ElementType.PARAGRAPH:
                body_parts.append(str(elem.content))
                continue

        # Add body text to content placeholder
        if body_parts:
            # Find the body placeholder (index 1 typically)
            body_ph = None
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == 1:
                    body_ph = ph
                    break
            if body_ph and hasattr(body_ph, "text_frame"):
                tf = body_ph.text_frame
                tf.clear()
                for i, text in enumerate(body_parts):
                    if i == 0:
                        tf.paragraphs[0].text = text
                    else:
                        tf.add_paragraph().text = text

        # Speaker notes
        if notes_text:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes_text

    def _add_table_to_slide(self, slide: Any, rows: list[list[str]]) -> None:
        from pptx.util import Inches, Emu

        if not rows:
            return

        n_rows = len(rows)
        n_cols = len(rows[0]) if rows[0] else 1

        # Position table in the middle area of the slide
        left = Inches(0.5)
        top = Inches(2.0)
        width = Inches(9.0)
        height = Inches(0.4 * n_rows)

        table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
        table = table_shape.table

        for ri, row in enumerate(rows):
            for ci, cell_text in enumerate(row):
                if ci < len(table.columns):
                    table.cell(ri, ci).text = str(cell_text)

    def _add_image_to_slide(self, slide: Any, elem: Element, model: DocumentModel) -> None:
        from pptx.util import Inches

        img_name = Path(str(elem.content)).name
        if img_name not in model.images:
            return

        img_data = model.images[img_name]
        img_stream = io.BytesIO(img_data.data)

        try:
            left = Inches(1.0)
            top = Inches(2.0)
            width = Inches(5.0)
            slide.shapes.add_picture(img_stream, left, top, width=width)
        except Exception as exc:
            log.debug("pptx.add_image_failed", error=str(exc))

    def _try_tier3_export(
        self,
        model: DocumentModel,
        output_path: Path,
        original_path: Path,
        sidecar: dict[str, Any],
    ) -> bool:
        """Attempt Tier 3: patch text in original PPTX. Returns True on success."""
        from pptx import Presentation

        try:
            prs = Presentation(str(original_path))
        except Exception:
            return False

        # Build a map of content hashes from the model
        model_texts: dict[int, list[str]] = {}  # slide_idx → list of texts
        slide_idx = -1
        for elem in model.elements:
            if elem.type == ElementType.HEADING and elem.level == 2:
                slide_idx += 1
                model_texts.setdefault(slide_idx, [])
                model_texts[slide_idx].append(str(elem.content))
            elif elem.type == ElementType.PARAGRAPH:
                model_texts.setdefault(max(slide_idx, 0), []).append(str(elem.content))

        # Check hash match — need ≥80% of original text still present
        original_texts: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            original_texts.append(para.text.strip())

        if not original_texts:
            return False

        model_text_set = set()
        for texts in model_texts.values():
            model_text_set.update(texts)

        match_count = sum(1 for t in original_texts if t in model_text_set)
        match_ratio = match_count / len(original_texts)

        if match_ratio < 0.8:
            return False

        # Patch: update text in slides
        for si, slide in enumerate(prs.slides):
            if si not in model_texts:
                continue
            text_idx = 0
            texts = model_texts[si]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip() and text_idx < len(texts):
                            # Preserve formatting of first run, update text
                            if para.runs:
                                para.runs[0].text = texts[text_idx]
                                for run in para.runs[1:]:
                                    run.text = ""
                            else:
                                para.text = texts[text_idx]
                            text_idx += 1

        prs.save(str(output_path))
        return True

    # ── Style extraction ──────────────────────────────────────────────────────

    def extract_styles(self, file_path: Path) -> dict[str, Any]:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        file_path = Path(file_path)
        _tmp_pptx: Path | None = None
        if file_path.suffix.lower() == ".ppt":
            from core.libreoffice_helper import convert_with_libreoffice

            file_path = convert_with_libreoffice(file_path, "pptx")
            _tmp_pptx = file_path

        try:
            return self._extract_styles_impl(file_path, Presentation, MSO_SHAPE_TYPE)
        finally:
            if _tmp_pptx and _tmp_pptx.exists():
                _tmp_pptx.unlink(missing_ok=True)

    def _extract_styles_impl(self, file_path: Path, Presentation: Any, MSO_SHAPE_TYPE: Any) -> dict[str, Any]:
        from collections import Counter
        styles: dict[str, Any] = {"document_level": {}}
        _hash_counter: Counter[str] = Counter()

        prs = Presentation(str(file_path))
        styles["document_level"]["slide_width"] = prs.slide_width
        styles["document_level"]["slide_height"] = prs.slide_height

        for idx, slide in enumerate(prs.slides):
            slide_key = f"slide_{idx + 1}"
            slide_style: dict[str, Any] = {
                "layout_name": slide.slide_layout.name,
                "layout_index": list(prs.slide_layouts).index(slide.slide_layout)
                if slide.slide_layout in prs.slide_layouts
                else 1,
            }

            shapes_info: list[dict] = []
            for shape in slide.shapes:
                shape_info: dict[str, Any] = {
                    "name": shape.name,
                    "left": shape.left,
                    "top": shape.top,
                    "width": shape.width,
                    "height": shape.height,
                }
                try:
                    if shape.placeholder_format is not None:
                        shape_info["placeholder_idx"] = shape.placeholder_format.idx
                except (ValueError, AttributeError):
                    pass

                # Font info from text runs
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font:
                                font_info: dict[str, Any] = {}
                                if run.font.name:
                                    font_info["name"] = run.font.name
                                if run.font.size:
                                    font_info["size"] = run.font.size
                                if run.font.bold is not None:
                                    font_info["bold"] = run.font.bold
                                if run.font.italic is not None:
                                    font_info["italic"] = run.font.italic
                                try:
                                    if run.font.color and run.font.color.rgb:
                                        font_info["color"] = str(run.font.color.rgb)
                                except (AttributeError, TypeError):
                                    pass
                                if font_info:
                                    shape_info["font"] = font_info
                                    break
                        break

                shapes_info.append(shape_info)

            slide_style["shapes"] = shapes_info
            styles[slide_key] = slide_style

            # Also key by content hash for sidecar lookup
            title = self._get_slide_title(slide, idx + 1)
            h = compute_content_hash(title)
            n = _hash_counter[h]
            _hash_counter[h] += 1
            styles[f"{h}:{n}"] = {"layout_index": slide_style.get("layout_index", 1)}

        return styles

    @classmethod
    def supports_format(cls, extension: str) -> bool:
        return extension.lower().lstrip(".") in cls.EXTENSIONS
