"""
Shared fixtures for MarkFlow test suite.

Provides: async test client, temporary directory fixtures, SQLite DB fixtures,
and runs generate_fixtures.py to ensure test files exist before any test runs.
"""

import os
import tempfile
from pathlib import Path

import pytest
import pytest_asyncio
from httpx import ASGITransport, AsyncClient

# ── Set env vars before any app module is imported ───────────────────────────
# aiosqlite :memory: creates a new DB per connection and won't persist across
# requests. Use a real temp file for the test session instead.
_TEST_DB = Path(tempfile.mktemp(suffix="_markflow_test.db"))
os.environ["DB_PATH"] = str(_TEST_DB)

# Enable auth bypass so all existing tests work without credentials
os.environ["DEV_BYPASS_AUTH"] = "true"

# Generate test DOCX fixtures before anything runs
from tests.generate_fixtures import generate_all  # noqa: E402

generate_all()


# ── App + async fixtures ──────────────────────────────────────────────────────

@pytest.fixture(scope="session")
def anyio_backend():
    return "asyncio"


@pytest_asyncio.fixture(scope="session")
async def client():
    """
    Async HTTP client pointed at the FastAPI app with a session-scoped temp DB.

    We explicitly call init_db() because ASGITransport does not guarantee the
    FastAPI lifespan runs before the first test request.
    """
    from core.database import init_db
    from main import app

    # Ensure schema is created before any test uses the DB
    await init_db()

    transport = ASGITransport(app=app)
    async with AsyncClient(transport=transport, base_url="http://test") as ac:
        yield ac

    # Clean up test DB after session
    _TEST_DB.unlink(missing_ok=True)


# ── Directory fixtures ────────────────────────────────────────────────────────

@pytest.fixture
def fixtures_dir() -> Path:
    """Path to tests/fixtures/."""
    return Path(__file__).parent / "fixtures"


@pytest.fixture
def simple_docx(fixtures_dir) -> Path:
    return fixtures_dir / "simple.docx"


@pytest.fixture
def complex_docx(fixtures_dir) -> Path:
    return fixtures_dir / "complex.docx"


# ── PDF fixtures ─────────────────────────────────────────────────────────────

@pytest.fixture
def simple_text_pdf(fixtures_dir) -> Path:
    return fixtures_dir / "simple_text.pdf"


@pytest.fixture
def scanned_pdf(fixtures_dir) -> Path:
    return fixtures_dir / "scanned.pdf"


@pytest.fixture
def mixed_pdf(fixtures_dir) -> Path:
    return fixtures_dir / "mixed.pdf"


# ── PPTX fixtures ────────────────────────────────────────────────────────────

@pytest.fixture
def simple_pptx(fixtures_dir) -> Path:
    return fixtures_dir / "simple.pptx"


# ── XLSX fixtures ────────────────────────────────────────────────────────────

@pytest.fixture
def simple_xlsx(fixtures_dir) -> Path:
    return fixtures_dir / "simple.xlsx"


@pytest.fixture
def complex_xlsx(fixtures_dir) -> Path:
    return fixtures_dir / "complex.xlsx"


# ── CSV/TSV fixtures ─────────────────────────────────────────────────────────

@pytest.fixture
def simple_csv(fixtures_dir) -> Path:
    return fixtures_dir / "simple.csv"


@pytest.fixture
def unicode_csv(fixtures_dir) -> Path:
    return fixtures_dir / "unicode.csv"


@pytest.fixture
def simple_tsv(fixtures_dir) -> Path:
    return fixtures_dir / "simple.tsv"


# ── Additional fixtures for Phase 5 ────────────────────────────────────────

@pytest.fixture
def latin1_csv_path(tmp_path) -> Path:
    """CSV encoded in latin-1 to test encoding detection."""
    import csv

    rows = [
        ["id", "name", "city"],
        ["1", "José", "São Paulo"],
        ["2", "François", "Genève"],
        ["3", "Günter", "Zürich"],
    ]
    path = tmp_path / "latin1.csv"
    with open(path, "w", newline="", encoding="latin-1") as f:
        writer = csv.writer(f)
        writer.writerows(rows)
    return path


@pytest.fixture
def simple_pdf_path(tmp_path) -> Path:
    """Generates a minimal text-native PDF using fpdf2."""
    from fpdf import FPDF

    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 16)
    pdf.cell(0, 12, "Test PDF Title", new_x="LMARGIN", new_y="NEXT")
    pdf.set_font("Helvetica", "", 12)
    pdf.multi_cell(0, 7,
        "This is a test paragraph for the PDF handler. "
        "It contains enough text to pass the minimum text length threshold."
    )
    path = tmp_path / "test.pdf"
    pdf.output(str(path))
    return path


@pytest.fixture
def simple_pptx_path(tmp_path) -> Path:
    """3-slide PPTX: title slide, content slide with bullet list, slide with table."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    # Slide 1 — Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Test Title"
    slide.placeholders[1].text = "Subtitle text"

    # Slide 2 — Content with bullets
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Content Slide"
    body = slide2.placeholders[1]
    body.text = "Bullet point one"
    body.text_frame.add_paragraph().text = "Bullet point two"

    # Slide 3 — Table
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Table Slide"
    tbl = slide3.shapes.add_table(3, 2, Inches(1), Inches(2), Inches(6), Inches(2)).table
    tbl.cell(0, 0).text = "Name"
    tbl.cell(0, 1).text = "Value"
    tbl.cell(1, 0).text = "Alpha"
    tbl.cell(1, 1).text = "100"
    tbl.cell(2, 0).text = "Beta"
    tbl.cell(2, 1).text = "200"

    path = tmp_path / "test.pptx"
    prs.save(str(path))
    return path


@pytest.fixture
def simple_xlsx_path(tmp_path) -> Path:
    """Workbook: Sheet1 = simple table, Sheet2 = merged cells, Sheet3 = formula cells."""
    import openpyxl

    wb = openpyxl.Workbook()
    ws1 = wb.active
    ws1.title = "Data"
    ws1.append(["Name", "Value"])
    ws1.append(["Alice", "100"])
    ws1.append(["Bob", "200"])
    ws1.append(["Charlie", "300"])

    ws2 = wb.create_sheet("Merged")
    ws2.append(["Category", "Count"])
    ws2.append(["A", "5"])
    ws2.append(["A", "3"])
    ws2.merge_cells("A2:A3")

    ws3 = wb.create_sheet("Formulas")
    ws3.append(["X", "Y", "Sum"])
    ws3.append([10, 20, "=A2+B2"])
    ws3.append([30, 40, "=A3+B3"])

    path = tmp_path / "test.xlsx"
    wb.save(str(path))
    return path


@pytest.fixture
def simple_csv_path(tmp_path) -> Path:
    """UTF-8 CSV, 5 columns, 20 rows. Header row."""
    import csv

    path = tmp_path / "test.csv"
    with open(path, "w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f)
        writer.writerow(["id", "name", "age", "city", "score"])
        for i in range(1, 21):
            writer.writerow([str(i), f"Person_{i}", str(20 + i), f"City_{i}", f"{50.0 + i}"])
    return path


@pytest.fixture
def tsv_path(tmp_path) -> Path:
    """Tab-separated file."""
    import csv

    path = tmp_path / "test.tsv"
    with open(path, "w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f, delimiter="\t")
        writer.writerow(["id", "name", "value"])
        for i in range(1, 6):
            writer.writerow([str(i), f"Item_{i}", str(i * 100)])
    return path


@pytest.fixture
def document_model_with_all_elements():
    """DocumentModel containing every ElementType — used to test round-trip completeness."""
    from core.document_model import DocumentModel, DocumentMetadata, Element, ElementType

    model = DocumentModel()
    model.metadata = DocumentMetadata(
        source_file="all_elements.md",
        source_format="md",
        title="All Elements Test",
    )
    model.add_element(Element(type=ElementType.HEADING, content="Main Title", level=1))
    model.add_element(Element(type=ElementType.HEADING, content="Section", level=2))
    model.add_element(Element(type=ElementType.PARAGRAPH, content="A regular paragraph."))
    model.add_element(Element(type=ElementType.CODE_BLOCK, content="print('hello')"))
    model.add_element(Element(type=ElementType.BLOCKQUOTE, content="A wise quote."))
    model.add_element(Element(type=ElementType.HORIZONTAL_RULE, content=""))
    model.add_element(Element(type=ElementType.LIST_ITEM, content="First item", attributes={"ordered": False}))
    model.add_element(Element(type=ElementType.LIST_ITEM, content="Second item", attributes={"ordered": True}))
    model.add_element(Element(
        type=ElementType.TABLE,
        content=[["H1", "H2"], ["A", "B"], ["C", "D"]],
    ))
    model.add_element(Element(type=ElementType.PAGE_BREAK, content=""))
    model.add_element(Element(type=ElementType.PARAGRAPH, content="After page break."))
    model.add_element(Element(
        type=ElementType.IMAGE,
        content="assets/test.png",
        attributes={"alt": "test image", "src": "assets/test.png"},
    ))
    model.add_element(Element(
        type=ElementType.FOOTNOTE,
        content="This is a footnote.",
        attributes={"id": "1"},
    ))
    return model


# ── Active Operations Registry test fixtures (v0.35.0) ──────────────────────

@pytest_asyncio.fixture(autouse=True)
async def _set_hydration_event():
    """Tests don't run lifespan hydration — fake the event so register_op
    doesn't block. Also register a stub cancel hook for the op_types that
    Task 3's tests exercise as cancellable=True. Real hooks are registered
    by their owning subsystems at module import (Tasks 14, 16, 17, 18, 19,
    20, 23); until those land, tests that need cancellable=True for a
    given op_type must wire it up here.

    Caveat: this fixture leaves `_cancel_hooks["pipeline.run_now"]`
    populated for the entire suite (no teardown — `setdefault` is
    idempotent). Any future test that needs to assert "no real hook
    is registered for pipeline.run_now" must explicitly
    `_cancel_hooks.pop("pipeline.run_now", None)` inside the test
    body — or this fixture should be reshaped to yield-then-clear."""
    from core import active_ops
    active_ops._hydration_complete.set()

    # Stub hook for pipeline.run_now so test_register_op_returns_unique_uuid_and_persists
    # (which passes cancellable=True) succeeds. Intentionally NOT registering
    # for db.backup so test_register_op_cancellable_without_hook_raises still raises.
    async def _stub_cancel(op_id: str) -> None:
        return None

    active_ops._cancel_hooks.setdefault("pipeline.run_now", _stub_cancel)

    yield
    # Don't clear — once set, leave set for the rest of the suite


# ── authed_operator / authed_manager: role-scoped AsyncClients (v0.35.0) ─────
#
# Both fixtures inject a concrete role via dependency_overrides[get_current_user].
# When a test uses BOTH fixtures simultaneously (e.g. test_cancel_requires_manager_role),
# the overrides would conflict if stored globally. We solve this with a ContextVar:
# the override function reads the desired role from a ContextVar, and each client
# sets the ContextVar before each request via an httpx event hook. This ensures
# role isolation even when both clients are live concurrently.

import contextvars as _cv
_test_user_role: "_cv.ContextVar[object | None]" = _cv.ContextVar(
    "_test_user_role", default=None
)


def _make_role_scoped_client(app, role):
    """Build an AsyncClient that injects `role` into every request via ContextVar."""
    from core.auth import AuthenticatedUser, UserRole, get_current_user

    # Install (or replace) the override that reads the ContextVar
    async def _role_from_contextvar():
        _role = _test_user_role.get()
        if _role is None:
            _role = role  # fallback for single-fixture tests
        return AuthenticatedUser(
            sub=f"test-{_role.value}",
            email=f"{_role.value}@test",
            role=_role,
            is_service_account=False,
        )

    app.dependency_overrides[get_current_user] = _role_from_contextvar

    # httpx event hook: set the ContextVar to this fixture's role before each request
    async def _set_role(request):  # noqa: ARG001
        _test_user_role.set(role)

    transport = ASGITransport(app=app)
    return AsyncClient(
        transport=transport,
        base_url="http://test",
        event_hooks={"request": [_set_role]},
    )


@pytest_asyncio.fixture
async def authed_operator(client):
    """Per-test AsyncClient with exactly OPERATOR role.

    Injects UserRole.OPERATOR via a ContextVar-backed dependency override so
    that endpoints guarded by require_role(UserRole.MANAGER) return 403, while
    endpoints guarded by require_role(UserRole.OPERATOR) return 200.
    Safe to use alongside authed_manager in the same test.

    Depends on the session-scoped ``client`` fixture only to ensure
    ``init_db()`` has run (which applies all _MIGRATIONS, including v29's
    ``active_operations`` table) before this fixture yields. Without this
    dependency, running endpoint tests in isolation skips migrations and
    DB-touching tests fail with ``no such table``."""
    from core.auth import UserRole, get_current_user
    from main import app

    ac = _make_role_scoped_client(app, UserRole.OPERATOR)
    try:
        async with ac:
            yield ac
    finally:
        app.dependency_overrides.pop(get_current_user, None)


@pytest_asyncio.fixture
async def authed_manager(client):
    """Per-test AsyncClient with exactly MANAGER role.

    Injects UserRole.MANAGER via a ContextVar-backed dependency override so
    that endpoints guarded by require_role(UserRole.MANAGER) return 200.
    Safe to use alongside authed_operator in the same test.

    Depends on the session-scoped ``client`` fixture only to ensure
    ``init_db()`` has run (which applies all _MIGRATIONS, including v29's
    ``active_operations`` table) before this fixture yields. Without this
    dependency, running endpoint tests in isolation skips migrations and
    DB-touching tests fail with ``no such table``."""
    from core.auth import UserRole, get_current_user
    from main import app

    ac = _make_role_scoped_client(app, UserRole.MANAGER)
    try:
        async with ac:
            yield ac
    finally:
        app.dependency_overrides.pop(get_current_user, None)
