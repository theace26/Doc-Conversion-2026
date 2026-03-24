"""Tests for formats/pptx_handler.py — PPTX ingest, export, style extraction, round-trip."""

import tempfile
from pathlib import Path

import pytest

from core.document_model import ElementType
from formats.pptx_handler import PptxHandler


@pytest.fixture
def handler():
    return PptxHandler()


# ── Ingest: simple.pptx ─────────────────────────────────────────────────────

def test_ingest_returns_model(simple_pptx, handler):
    model = handler.ingest(simple_pptx)
    assert model is not None
    assert len(model.elements) > 0


def test_ingest_slide_count(simple_pptx, handler):
    """5 slides should produce 5 H2 headings."""
    model = handler.ingest(simple_pptx)
    headings = model.get_elements_by_type(ElementType.HEADING)
    h2s = [h for h in headings if h.level == 2]
    assert len(h2s) == 5


def test_ingest_slide_titles(simple_pptx, handler):
    model = handler.ingest(simple_pptx)
    headings = model.get_elements_by_type(ElementType.HEADING)
    titles = [h.content for h in headings if h.level == 2]
    assert "Test Presentation" in titles
    assert "Introduction" in titles
    assert "Data Overview" in titles
    assert "Image Slide" in titles
    assert "Summary" in titles


def test_ingest_body_text(simple_pptx, handler):
    model = handler.ingest(simple_pptx)
    paras = model.get_elements_by_type(ElementType.PARAGRAPH)
    para_texts = [str(p.content) for p in paras]
    assert any("first paragraph" in t for t in para_texts)


def test_ingest_has_table(simple_pptx, handler):
    model = handler.ingest(simple_pptx)
    tables = model.get_elements_by_type(ElementType.TABLE)
    assert len(tables) >= 1


def test_ingest_table_dimensions(simple_pptx, handler):
    model = handler.ingest(simple_pptx)
    tables = model.get_elements_by_type(ElementType.TABLE)
    assert len(tables) >= 1
    rows = tables[0].content
    assert len(rows) == 4  # 1 header + 3 data
    assert len(rows[0]) == 3  # 3 columns


def test_ingest_table_content(simple_pptx, handler):
    model = handler.ingest(simple_pptx)
    tables = model.get_elements_by_type(ElementType.TABLE)
    rows = tables[0].content
    assert "Name" in rows[0]
    assert "Alice" in rows[1]


def test_ingest_has_image(simple_pptx, handler):
    model = handler.ingest(simple_pptx)
    images = model.get_elements_by_type(ElementType.IMAGE)
    assert len(images) >= 1


def test_ingest_image_data_extracted(simple_pptx, handler):
    model = handler.ingest(simple_pptx)
    assert len(model.images) >= 1
    for name, img_data in model.images.items():
        assert len(img_data.data) > 0


def test_ingest_speaker_notes(simple_pptx, handler):
    model = handler.ingest(simple_pptx)
    blockquotes = model.get_elements_by_type(ElementType.BLOCKQUOTE)
    notes = [bq for bq in blockquotes if "Speaker Notes:" in str(bq.content)]
    assert len(notes) >= 1
    assert "key findings" in str(notes[0].content)


def test_ingest_has_horizontal_rules(simple_pptx, handler):
    """Horizontal rules separate slides."""
    model = handler.ingest(simple_pptx)
    rules = model.get_elements_by_type(ElementType.HORIZONTAL_RULE)
    assert len(rules) == 4  # 5 slides → 4 separators


def test_ingest_metadata(simple_pptx, handler):
    model = handler.ingest(simple_pptx)
    assert model.metadata.source_format == "pptx"
    assert model.metadata.page_count == 5


def test_ingest_bold_text_preserved(simple_pptx, handler):
    """Bold formatting should be preserved as markdown markers."""
    model = handler.ingest(simple_pptx)
    paras = model.get_elements_by_type(ElementType.PARAGRAPH)
    texts = [str(p.content) for p in paras]
    assert any("**" in t for t in texts)


# ── Export ───────────────────────────────────────────────────────────────────

def test_export_creates_file(simple_pptx, handler, tmp_path):
    model = handler.ingest(simple_pptx)
    output = tmp_path / "output.pptx"
    handler.export(model, output)
    assert output.exists()
    assert output.stat().st_size > 0


def test_export_slide_count(simple_pptx, handler, tmp_path):
    from pptx import Presentation

    model = handler.ingest(simple_pptx)
    output = tmp_path / "output.pptx"
    handler.export(model, output)

    prs = Presentation(str(output))
    assert len(prs.slides) == 5


def test_export_slide_titles(simple_pptx, handler, tmp_path):
    from pptx import Presentation

    model = handler.ingest(simple_pptx)
    output = tmp_path / "output.pptx"
    handler.export(model, output)

    prs = Presentation(str(output))
    titles = []
    for slide in prs.slides:
        if slide.shapes.title:
            titles.append(slide.shapes.title.text)
    assert "Test Presentation" in titles
    assert "Introduction" in titles


def test_export_has_table(simple_pptx, handler, tmp_path):
    from pptx import Presentation

    model = handler.ingest(simple_pptx)
    output = tmp_path / "output.pptx"
    handler.export(model, output)

    prs = Presentation(str(output))
    table_found = False
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table:
                table_found = True
                assert len(shape.table.rows) == 4
                assert len(shape.table.columns) == 3
    assert table_found


def test_export_speaker_notes_restored(simple_pptx, handler, tmp_path):
    from pptx import Presentation

    model = handler.ingest(simple_pptx)
    output = tmp_path / "output.pptx"
    handler.export(model, output)

    prs = Presentation(str(output))
    notes_found = False
    for slide in prs.slides:
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text
            if "key findings" in notes_text:
                notes_found = True
    assert notes_found


# ── Style extraction ─────────────────────────────────────────────────────────

def test_extract_styles_document_level(simple_pptx, handler):
    styles = handler.extract_styles(simple_pptx)
    assert "document_level" in styles
    assert "slide_width" in styles["document_level"]
    assert "slide_height" in styles["document_level"]


def test_extract_styles_has_slide_entries(simple_pptx, handler):
    styles = handler.extract_styles(simple_pptx)
    assert "slide_1" in styles
    assert "slide_2" in styles


def test_extract_styles_layout_info(simple_pptx, handler):
    styles = handler.extract_styles(simple_pptx)
    slide1 = styles.get("slide_1", {})
    assert "layout_name" in slide1


# ── Round-trip ───────────────────────────────────────────────────────────────

def test_roundtrip_slide_count(simple_pptx, handler, tmp_path):
    """PPTX → DocumentModel → PPTX: 5 slides survive."""
    from pptx import Presentation

    model = handler.ingest(simple_pptx)
    output = tmp_path / "roundtrip.pptx"
    handler.export(model, output)

    prs = Presentation(str(output))
    assert len(prs.slides) == 5


def test_roundtrip_titles_match(simple_pptx, handler, tmp_path):
    from pptx import Presentation

    model = handler.ingest(simple_pptx)
    output = tmp_path / "roundtrip.pptx"
    handler.export(model, output)

    prs = Presentation(str(output))
    titles = [
        slide.shapes.title.text
        for slide in prs.slides
        if slide.shapes.title
    ]
    assert "Test Presentation" in titles
    assert "Summary" in titles


def test_roundtrip_content_present(simple_pptx, handler, tmp_path):
    """Re-ingest exported PPTX and check content."""
    model = handler.ingest(simple_pptx)
    output = tmp_path / "roundtrip.pptx"
    handler.export(model, output)

    model2 = handler.ingest(output)
    h2s = [h for h in model2.get_elements_by_type(ElementType.HEADING) if h.level == 2]
    assert len(h2s) == 5


# ── supports_format ──────────────────────────────────────────────────────────

def test_supports_pptx():
    assert PptxHandler.supports_format("pptx")
    assert PptxHandler.supports_format(".pptx")
    assert not PptxHandler.supports_format("xlsx")
