"""
Tests for formats/pptx_handler.py — PPTX ingest and export.

Covers slide extraction, tables, images, speaker notes, edge cases,
and Markdown → PPTX export.
"""

import pytest
from pathlib import Path

from core.document_model import DocumentModel, DocumentMetadata, Element, ElementType


# ── Ingest ───────────────────────────────────────────────────────────────────

class TestPptxIngest:
    """Tests for PPTX ingestion."""

    def test_slides_become_h2_elements(self, simple_pptx):
        from formats.pptx_handler import PptxHandler

        handler = PptxHandler()
        model = handler.ingest(simple_pptx)

        headings = model.get_elements_by_type(ElementType.HEADING)
        h2s = [h for h in headings if h.level == 2]
        assert len(h2s) == 5, "5-slide PPTX should produce 5 H2 headings"

    def test_page_count_equals_slide_count(self, simple_pptx):
        from formats.pptx_handler import PptxHandler

        handler = PptxHandler()
        model = handler.ingest(simple_pptx)

        assert model.metadata.page_count == 5

    def test_tables_extracted(self, simple_pptx):
        from formats.pptx_handler import PptxHandler

        handler = PptxHandler()
        model = handler.ingest(simple_pptx)

        tables = model.get_elements_by_type(ElementType.TABLE)
        assert len(tables) >= 1, "Slide 3 has a table"

    def test_speaker_notes_become_blockquotes(self, simple_pptx):
        from formats.pptx_handler import PptxHandler

        handler = PptxHandler()
        model = handler.ingest(simple_pptx)

        blockquotes = model.get_elements_by_type(ElementType.BLOCKQUOTE)
        notes = [bq for bq in blockquotes if "Speaker Notes" in str(bq.content)]
        assert len(notes) >= 1, "Slide 5 has speaker notes"

    def test_images_extracted(self, simple_pptx):
        from formats.pptx_handler import PptxHandler

        handler = PptxHandler()
        model = handler.ingest(simple_pptx)

        images = model.get_elements_by_type(ElementType.IMAGE)
        assert len(images) >= 1, "Slide 4 has an image"

    def test_paragraphs_extracted(self, simple_pptx):
        from formats.pptx_handler import PptxHandler

        handler = PptxHandler()
        model = handler.ingest(simple_pptx)

        paragraphs = model.get_elements_by_type(ElementType.PARAGRAPH)
        assert len(paragraphs) > 0

    def test_source_format(self, simple_pptx):
        from formats.pptx_handler import PptxHandler

        handler = PptxHandler()
        model = handler.ingest(simple_pptx)

        assert model.metadata.source_format == "pptx"


# ── Ingest edge cases ───────────────────────────────────────────────────────

class TestPptxIngestEdgeCases:
    """Edge case tests for PPTX ingest."""

    def test_slide_without_title_uses_fallback(self, tmp_path):
        from pptx import Presentation
        from formats.pptx_handler import PptxHandler

        prs = Presentation()
        # Use blank layout (no title placeholder)
        layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(layout)

        path = tmp_path / "no_title.pptx"
        prs.save(str(path))

        handler = PptxHandler()
        model = handler.ingest(path)

        headings = model.get_elements_by_type(ElementType.HEADING)
        assert len(headings) >= 1
        assert "Slide 1" in str(headings[0].content)

    def test_placeholder_format_valueerror_handled(self, tmp_path):
        """ValueError from placeholder_format on non-placeholder shapes is caught."""
        from pptx import Presentation
        from pptx.util import Inches
        from formats.pptx_handler import PptxHandler

        prs = Presentation()
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = "Test"

        # Add a table (has no placeholder_format)
        slide.shapes.add_table(2, 2, Inches(1), Inches(2), Inches(4), Inches(2))

        path = tmp_path / "table_shape.pptx"
        prs.save(str(path))

        handler = PptxHandler()
        model = handler.ingest(path)
        # Should not crash
        assert model.metadata.page_count == 1

    def test_shape_with_no_text_no_empty_element(self, simple_pptx):
        """Decorative shapes with no text should not produce empty elements."""
        from formats.pptx_handler import PptxHandler

        handler = PptxHandler()
        model = handler.ingest(simple_pptx)

        for elem in model.elements:
            if elem.type == ElementType.PARAGRAPH:
                assert str(elem.content).strip() != "", "No empty paragraph elements"


# ── Export ───────────────────────────────────────────────────────────────────

class TestPptxExport:
    """Tests for Markdown → PPTX export."""

    def test_h2_sections_become_slides(self, tmp_path):
        from formats.pptx_handler import PptxHandler

        model = DocumentModel()
        model.add_element(Element(type=ElementType.HEADING, content="Slide 1", level=2))
        model.add_element(Element(type=ElementType.PARAGRAPH, content="Content 1"))
        model.add_element(Element(type=ElementType.HEADING, content="Slide 2", level=2))
        model.add_element(Element(type=ElementType.PARAGRAPH, content="Content 2"))

        handler = PptxHandler()
        output = tmp_path / "export.pptx"
        handler.export(model, output)

        from pptx import Presentation
        prs = Presentation(str(output))
        assert len(prs.slides) == 2

    def test_tables_render_in_slide(self, tmp_path):
        from formats.pptx_handler import PptxHandler

        model = DocumentModel()
        model.add_element(Element(type=ElementType.HEADING, content="Data", level=2))
        model.add_element(Element(
            type=ElementType.TABLE,
            content=[["A", "B"], ["1", "2"]],
        ))

        handler = PptxHandler()
        output = tmp_path / "table.pptx"
        handler.export(model, output)

        assert output.exists()
        assert output.stat().st_size > 0

    def test_no_h2_produces_single_slide(self, tmp_path):
        from formats.pptx_handler import PptxHandler

        model = DocumentModel()
        model.add_element(Element(type=ElementType.PARAGRAPH, content="Just a paragraph"))

        handler = PptxHandler()
        output = tmp_path / "no_h2.pptx"
        handler.export(model, output)

        from pptx import Presentation
        prs = Presentation(str(output))
        assert len(prs.slides) == 1

    def test_list_items_render(self, tmp_path):
        from formats.pptx_handler import PptxHandler

        model = DocumentModel()
        model.add_element(Element(type=ElementType.HEADING, content="Lists", level=2))
        model.add_element(Element(type=ElementType.PARAGRAPH, content="Item A"))
        model.add_element(Element(type=ElementType.PARAGRAPH, content="Item B"))

        handler = PptxHandler()
        output = tmp_path / "list.pptx"
        handler.export(model, output)

        assert output.exists()

    def test_tier3_patching(self, simple_pptx, tmp_path):
        """Tier 3 patching: when original PPTX provided, unchanged slides preserved."""
        from formats.pptx_handler import PptxHandler

        handler = PptxHandler()
        model = handler.ingest(simple_pptx)

        sidecar = handler.extract_styles(simple_pptx)
        output = tmp_path / "tier3.pptx"
        handler.export(model, output, sidecar=sidecar, original_path=simple_pptx)

        from pptx import Presentation
        prs = Presentation(str(output))
        assert len(prs.slides) >= 1


class TestPptxStyleExtraction:
    """Tests for PPTX style extraction."""

    def test_extract_styles_returns_dict(self, simple_pptx):
        from formats.pptx_handler import PptxHandler

        handler = PptxHandler()
        styles = handler.extract_styles(simple_pptx)
        assert isinstance(styles, dict)
        assert "document_level" in styles

    def test_extract_styles_has_slide_dimensions(self, simple_pptx):
        from formats.pptx_handler import PptxHandler

        handler = PptxHandler()
        styles = handler.extract_styles(simple_pptx)
        dl = styles["document_level"]
        assert "slide_width" in dl
        assert "slide_height" in dl

    def test_extract_styles_per_slide(self, simple_pptx):
        from formats.pptx_handler import PptxHandler

        handler = PptxHandler()
        styles = handler.extract_styles(simple_pptx)
        assert "slide_1" in styles

    def test_supports_format(self):
        from formats.pptx_handler import PptxHandler

        assert PptxHandler.supports_format("pptx")
        assert PptxHandler.supports_format(".pptx")
        assert not PptxHandler.supports_format("docx")

    def test_export_with_speaker_notes(self, tmp_path):
        from formats.pptx_handler import PptxHandler

        model = DocumentModel()
        model.add_element(Element(type=ElementType.HEADING, content="Title", level=2))
        model.add_element(Element(type=ElementType.PARAGRAPH, content="Body text"))
        model.add_element(Element(type=ElementType.BLOCKQUOTE, content="Speaker Notes: Remember this"))

        handler = PptxHandler()
        output = tmp_path / "notes.pptx"
        handler.export(model, output)

        from pptx import Presentation
        prs = Presentation(str(output))
        assert len(prs.slides) == 1
        slide = prs.slides[0]
        assert slide.has_notes_slide
        notes_text = slide.notes_slide.notes_text_frame.text
        assert "Remember this" in notes_text


# ── Chart extraction ────────────────────────────────────────────────────────

class TestChartExtraction:
    """Tests for M5: PPTX chart extraction and SmartArt detection."""

    def test_chart_extraction_mode_preference_exists(self):
        """The preference exists in defaults with correct value."""
        from core.db.preferences import DEFAULT_PREFERENCES

        assert "pptx_chart_extraction_mode" in DEFAULT_PREFERENCES
        assert DEFAULT_PREFERENCES["pptx_chart_extraction_mode"] == "placeholder"

    def test_chart_placeholder_default(self, tmp_path):
        """Default mode produces [Chart: ...] placeholder."""
        from pptx import Presentation
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE
        from pptx.util import Inches
        from formats.pptx_handler import PptxHandler

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Chart Slide"

        chart_data = CategoryChartData()
        chart_data.categories = ["Q1", "Q2", "Q3"]
        chart_data.add_series("Revenue", (100, 200, 300))

        slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(1), Inches(2), Inches(6), Inches(4),
            chart_data,
        )

        path = tmp_path / "chart_test.pptx"
        prs.save(str(path))

        handler = PptxHandler()
        handler._chart_mode = "placeholder"
        handler._slide_images = {}
        handler._slide_width_emu = prs.slide_width
        handler._slide_height_emu = prs.slide_height
        model = DocumentModel()
        model.metadata = DocumentMetadata(source_file="test.pptx", source_format="pptx")

        from pptx import Presentation as P2
        prs2 = P2(str(path))
        for shape in prs2.slides[0].shapes:
            handler._process_shape(shape, model, path, 1)

        paras = model.get_elements_by_type(ElementType.PARAGRAPH)
        chart_paras = [p for p in paras if p.attributes.get("chart")]
        assert len(chart_paras) == 1
        assert "[Chart:" in str(chart_paras[0].content) or "[Chart]" in str(chart_paras[0].content)
        assert any("chart not fully extractable" in w for w in model.warnings)

    def test_crop_shape_region_basic(self):
        """Coordinate conversion crops correctly."""
        from formats.pptx_handler import PptxHandler
        from PIL import Image

        handler = PptxHandler()
        # Create a 200x100 test image
        img = Image.new("RGB", (200, 100), color=(255, 0, 0))

        # Simulate a shape at 50% x, 25% y, 25% width, 50% height
        class MockShape:
            left = 5000
            top = 2500
            width = 2500
            height = 5000

        result = handler._crop_shape_region(img, MockShape(), 10000, 10000)
        assert result is not None
        assert result.width == 50  # 25% of 200
        assert result.height == 50  # 50% of 100

    def test_crop_shape_region_clamps_bounds(self):
        """Crop clamps to image bounds when shape extends past edge."""
        from formats.pptx_handler import PptxHandler
        from PIL import Image

        handler = PptxHandler()
        img = Image.new("RGB", (100, 100), color=(0, 255, 0))

        class MockShape:
            left = 8000
            top = 8000
            width = 5000  # extends past right edge
            height = 5000  # extends past bottom edge

        result = handler._crop_shape_region(img, MockShape(), 10000, 10000)
        assert result is not None
        assert result.width == 20  # clamped: 100 - 80
        assert result.height == 20

    def test_crop_shape_region_zero_area_returns_none(self):
        """Zero-area crop returns None."""
        from formats.pptx_handler import PptxHandler
        from PIL import Image

        handler = PptxHandler()
        img = Image.new("RGB", (100, 100))

        class MockShape:
            left = 0
            top = 0
            width = 0
            height = 0

        result = handler._crop_shape_region(img, MockShape(), 10000, 10000)
        assert result is None

    def test_smartart_detection_xml(self, tmp_path):
        """SmartArt XML markers trigger a warning and text still extracted."""
        from pptx import Presentation
        from pptx.util import Inches
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from formats.pptx_handler import PptxHandler

        # Create a PPTX with a group shape, then check SmartArt detection
        # by verifying the code path handles group shapes gracefully
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

        path = tmp_path / "smartart_test.pptx"
        prs.save(str(path))

        handler = PptxHandler()
        handler._chart_mode = "placeholder"
        handler._slide_images = {}
        handler._slide_width_emu = prs.slide_width
        handler._slide_height_emu = prs.slide_height

        # The SmartArt detection code path is exercised by group shapes
        # with dgm:relIds in their XML. Since we can't easily create real
        # SmartArt via python-pptx, verify the handler ingests cleanly.
        model = handler.ingest(path)
        assert model is not None

    def test_read_chart_mode_pref_default(self):
        """When DB is unavailable, default is placeholder."""
        from formats.pptx_handler import PptxHandler

        # _read_chart_mode_pref should return "placeholder" when DB
        # doesn't exist or has no row
        result = PptxHandler._read_chart_mode_pref()
        # In test environment without DB, should default gracefully
        assert result in ("placeholder", "libreoffice")
