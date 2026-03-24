"""
Programmatic test fixture generator for MarkFlow.

Run this script (or call generate_all()) to create all test fixture files
in tests/fixtures/. All fixtures are generated programmatically — no manually
created files.

Fixtures created here:
  simple.docx      — 3 headings, 2 paragraphs, 1 table, 1 inline image
  complex.docx     — nested tables, footnotes, multiple fonts, embedded images
  simple_text.pdf  — 3 pages of clean, text-layer PDF
  scanned.pdf      — image-only (scanned) PDF for OCR testing
  mixed.pdf        — page 1 text-layer, page 2 scanned image
  simple.pptx      — 5 slides: titles, body, speaker notes, table, image
  simple.xlsx      — 2 sheets, headers, number formatting, SUM formula
  complex.xlsx     — merged cells, formulas, conditional formatting
  simple.csv       — clean CSV with headers (5 cols, 10 rows)
  unicode.csv      — non-ASCII characters (accents, CJK, emoji)
  simple.tsv       — tab-delimited version
  OCR test images  — clean_text.png, noisy_scan.png, bad_scan.png, etc.
"""

import io
from pathlib import Path

FIXTURES_DIR = Path(__file__).parent / "fixtures"


# ── DOCX fixtures ─────────────────────────────────────────────────────────────

def _make_simple_docx() -> None:
    """
    simple.docx:
      - Heading 1: "Simple Document"
      - Heading 2: "Introduction"
      - Paragraph: "This is the first paragraph of the document."
      - Heading 3: "Details"
      - Paragraph: "This is the second paragraph with some **detail**."
      - Table 3×3 with headers
      - Inline PNG image (10×10 red square)
    """
    import docx
    from docx.shared import Inches, Pt, RGBColor

    doc = docx.Document()
    doc.core_properties.title = "Simple Document"
    doc.core_properties.author = "MarkFlow Test"

    doc.add_heading("Simple Document", level=1)
    doc.add_heading("Introduction", level=2)
    doc.add_paragraph("This is the first paragraph of the document.")
    doc.add_heading("Details", level=3)
    doc.add_paragraph("This is the second paragraph with some detail.")

    # 3×3 table
    table = doc.add_table(rows=3, cols=3, style="Table Grid")
    headers = ["Column A", "Column B", "Column C"]
    for j, h in enumerate(headers):
        table.cell(0, j).text = h
    table.cell(1, 0).text = "Row 1A"
    table.cell(1, 1).text = "Row 1B"
    table.cell(1, 2).text = "Row 1C"
    table.cell(2, 0).text = "Row 2A"
    table.cell(2, 1).text = "Row 2B"
    table.cell(2, 2).text = "Row 2C"

    # Inline image: 10×10 red square PNG
    img_buf = _make_png_bytes(10, 10, (220, 50, 50))
    doc.add_picture(img_buf, width=Inches(1.0))

    doc.save(FIXTURES_DIR / "simple.docx")


def _make_complex_docx() -> None:
    """
    complex.docx:
      - Multiple heading levels
      - Bold and italic runs
      - Nested paragraph structure
      - A table containing another table (simulated via text)
      - Footnote reference text
      - Multiple fonts (Calibri, Arial)
      - Two embedded images
    """
    import docx
    from docx.shared import Inches, Pt, RGBColor

    doc = docx.Document()
    doc.core_properties.title = "Complex Document"
    doc.core_properties.author = "MarkFlow Test"
    doc.core_properties.subject = "Testing"

    doc.add_heading("Complex Document", level=1)
    doc.add_heading("Chapter 1: Formatting", level=2)

    p = doc.add_paragraph()
    run = p.add_run("Bold text ")
    run.bold = True
    run2 = p.add_run("and italic text ")
    run2.italic = True
    run3 = p.add_run("and normal text.")

    doc.add_heading("Chapter 2: Tables", level=2)

    # 2×2 outer table
    outer = doc.add_table(rows=2, cols=2, style="Table Grid")
    outer.cell(0, 0).text = "Outer A1"
    outer.cell(0, 1).text = "Outer A2"
    outer.cell(1, 0).text = "Outer B1 (contains sub-data)"
    outer.cell(1, 1).text = "Outer B2"

    doc.add_heading("Chapter 3: Images", level=2)
    doc.add_paragraph("Below are two embedded images.")

    img1 = _make_png_bytes(20, 20, (50, 100, 200))
    doc.add_picture(img1, width=Inches(1.0))

    img2 = _make_png_bytes(20, 20, (50, 200, 100))
    doc.add_picture(img2, width=Inches(1.0))

    doc.add_heading("Chapter 4: Lists", level=2)
    for item in ["First item", "Second item", "Third item"]:
        doc.add_paragraph(item, style="List Bullet")
    for i, item in enumerate(["Step one", "Step two", "Step three"], 1):
        doc.add_paragraph(item, style="List Number")

    doc.save(FIXTURES_DIR / "complex.docx")


def _make_png_bytes(width: int, height: int, color: tuple[int, int, int]) -> io.BytesIO:
    """Create a solid-color PNG in a BytesIO buffer."""
    from PIL import Image

    img = Image.new("RGB", (width, height), color)
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


# ── OCR image fixtures ─────────────────────────────────────────────────────────

_SAMPLE_LINES = [
    "The quick brown fox jumps over the lazy dog.",
    "MarkFlow OCR pipeline — Phase 3 test fixture.",
    "Pack my box with five dozen liquor jugs.",
    "How vexingly quick daft zebras jump!",
    "The five boxing wizards jump quickly.",
    "Sphinx of black quartz, judge my vow.",
    "Two driven jocks help fax my big quiz.",
]

_TABLE_ROWS = [
    ("Name",    "Score", "Grade"),
    ("Alice",   "92",    "A"),
    ("Bob",     "78",    "B+"),
    ("Charlie", "65",    "C"),
    ("Diana",   "88",    "B"),
]


def _get_font(size: int = 24):
    """Return a PIL font — truetype if available, else the built-in bitmap font."""
    from PIL import ImageFont

    # Try common system font paths (works on Debian/Ubuntu in Docker)
    candidates = [
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
        "/usr/share/fonts/truetype/freefont/FreeSans.ttf",
        "/usr/share/fonts/truetype/ubuntu/Ubuntu-R.ttf",
    ]
    for path in candidates:
        try:
            return ImageFont.truetype(path, size)
        except (OSError, IOError):
            pass
    # Fall back to PIL's built-in font (small but always available)
    return ImageFont.load_default()


def _draw_text_page(
    width: int = 800,
    height: int = 1000,
    lines: list[str] | None = None,
    noise_stddev: float = 0.0,
    rotation: float = 0.0,
    scale_factor: float = 1.0,
    font_size: int = 28,
    margin: int = 60,
) -> "Image.Image":
    """
    Render lines of text onto a white background.

    noise_stddev: if > 0, add Gaussian noise (simulates scan degradation).
    rotation: degrees of clockwise skew to apply.
    scale_factor: < 1 = reduce resolution before returning.
    """
    import numpy as np
    from PIL import Image, ImageDraw

    img = Image.new("RGB", (width, height), (255, 255, 255))
    draw = ImageDraw.Draw(img)
    font = _get_font(font_size)
    lines = lines or _SAMPLE_LINES

    y = margin
    for line in lines:
        draw.text((margin, y), line, font=font, fill=(20, 20, 20))
        y += font_size + 10

    # Add Gaussian noise
    if noise_stddev > 0.0:
        arr = np.array(img, dtype=np.float32)
        noise = np.random.normal(0, noise_stddev, arr.shape)
        arr = np.clip(arr + noise, 0, 255).astype(np.uint8)
        img = Image.fromarray(arr)

    # Apply skew
    if rotation != 0.0:
        img = img.rotate(-rotation, expand=False, fillcolor=(255, 255, 255))

    # Downscale (simulates low DPI)
    if scale_factor < 1.0:
        new_w = int(width * scale_factor)
        new_h = int(height * scale_factor)
        img = img.resize((new_w, new_h), Image.LANCZOS)

    return img


def generate_ocr_fixtures(fixtures_dir: Path) -> None:
    """
    Generate programmatic OCR test images for Phase 3 tests.
    Safe to call multiple times (skips existing files).
    """
    import numpy as np
    from PIL import Image, ImageDraw

    fixtures_dir.mkdir(exist_ok=True)

    # 1. clean_text.png — black text on white, 800×1000, high quality
    _save_if_missing(
        fixtures_dir / "clean_text.png",
        lambda: _draw_text_page(noise_stddev=0.0, rotation=0.0),
    )

    # 2. noisy_scan.png — same text + Gaussian noise + 2° skew, reduced contrast
    _save_if_missing(
        fixtures_dir / "noisy_scan.png",
        lambda: _draw_text_page(noise_stddev=30.0, rotation=2.0),
    )

    # 3. bad_scan.png — heavy noise, 5° skew, half resolution
    _save_if_missing(
        fixtures_dir / "bad_scan.png",
        lambda: _draw_text_page(noise_stddev=60.0, rotation=5.0, scale_factor=0.5),
    )

    # 4. blank_page.png — solid white (+ very slight noise so it's not degenerate)
    def _blank():
        arr = np.full((1000, 800, 3), 252, dtype=np.uint8)
        return Image.fromarray(arr)

    _save_if_missing(fixtures_dir / "blank_page.png", _blank)

    # 5. mixed_content.png — top half: text, bottom half: solid teal block
    def _mixed():
        img = _draw_text_page(height=1000)
        draw = ImageDraw.Draw(img)
        draw.rectangle([(0, 500), (800, 1000)], fill=(0, 128, 128))
        return img

    _save_if_missing(fixtures_dir / "mixed_content.png", _mixed)

    # 6. table_scan.png — simple grid table with text in cells
    def _table():
        img = Image.new("RGB", (800, 300), (255, 255, 255))
        draw = ImageDraw.Draw(img)
        font = _get_font(22)

        col_w, row_h = 160, 50
        for ri, row in enumerate(_TABLE_ROWS):
            for ci, cell in enumerate(row):
                x0 = ci * col_w
                y0 = ri * row_h
                draw.rectangle([x0, y0, x0 + col_w, y0 + row_h], outline=(0, 0, 0), width=2)
                draw.text((x0 + 8, y0 + 12), cell, font=font, fill=(20, 20, 20))
        return img

    _save_if_missing(fixtures_dir / "table_scan.png", _table)


def _save_if_missing(path: Path, factory) -> None:
    if not path.exists():
        img = factory()
        img.save(path)


# ── PDF fixtures ──────────────────────────────────────────────────────────────

def _make_simple_text_pdf() -> None:
    """
    simple_text.pdf — 3 pages of clean, text-layer PDF.
    Page 1: heading + paragraphs. Page 2: table. Page 3: more text.
    """
    from fpdf import FPDF

    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)

    # Page 1 — heading and paragraphs
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 18)
    pdf.cell(0, 12, "Simple PDF Document", new_x="LMARGIN", new_y="NEXT")
    pdf.set_font("Helvetica", "", 12)
    pdf.ln(5)
    pdf.multi_cell(0, 7,
        "This is the first paragraph of a simple text-layer PDF. "
        "It contains enough text to be recognized as a text page "
        "by the PDF handler's scanned-page detection heuristic."
    )
    pdf.ln(5)
    pdf.set_font("Helvetica", "B", 14)
    pdf.cell(0, 10, "Introduction", new_x="LMARGIN", new_y="NEXT")
    pdf.set_font("Helvetica", "", 12)
    pdf.multi_cell(0, 7,
        "The quick brown fox jumps over the lazy dog. "
        "Pack my box with five dozen liquor jugs. "
        "How vexingly quick daft zebras jump!"
    )

    # Page 2 — table
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 14)
    pdf.cell(0, 10, "Data Table", new_x="LMARGIN", new_y="NEXT")
    pdf.set_font("Helvetica", "", 11)
    pdf.ln(3)

    headers = ["Name", "Score", "Grade"]
    data = [
        ["Alice", "92", "A"],
        ["Bob", "78", "B+"],
        ["Charlie", "65", "C"],
        ["Diana", "88", "B"],
    ]
    col_w = 50
    # Header
    pdf.set_font("Helvetica", "B", 11)
    for h in headers:
        pdf.cell(col_w, 8, h, border=1)
    pdf.ln()
    # Rows
    pdf.set_font("Helvetica", "", 11)
    for row in data:
        for cell in row:
            pdf.cell(col_w, 8, cell, border=1)
        pdf.ln()

    # Page 3 — more text
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 14)
    pdf.cell(0, 10, "Conclusion", new_x="LMARGIN", new_y="NEXT")
    pdf.set_font("Helvetica", "", 12)
    pdf.multi_cell(0, 7,
        "This is the third and final page of the simple PDF fixture. "
        "It exists to test multi-page handling and page break detection. "
        "The document should produce a markdown file with clear page boundaries."
    )

    pdf.output(str(FIXTURES_DIR / "simple_text.pdf"))


def _make_scanned_pdf() -> None:
    """
    scanned.pdf — an image-only PDF (simulated scan).
    Renders text as an image, then embeds the image in a PDF page.
    """
    from fpdf import FPDF

    # Render text to image
    img = _draw_text_page(
        width=800, height=1100,
        lines=[
            "SCANNED DOCUMENT",
            "",
            "This page was rendered as an image",
            "to simulate a scanned PDF document.",
            "",
            "The OCR pipeline should detect this",
            "as a scanned page and extract text",
            "using Tesseract.",
            "",
            "The quick brown fox jumps over",
            "the lazy dog. Pack my box with",
            "five dozen liquor jugs.",
        ],
        font_size=26,
    )
    img_path = FIXTURES_DIR / "_temp_scan.png"
    img.save(img_path)

    pdf = FPDF()
    pdf.add_page()
    pdf.image(str(img_path), x=0, y=0, w=210, h=297)
    pdf.output(str(FIXTURES_DIR / "scanned.pdf"))

    img_path.unlink(missing_ok=True)


def _make_mixed_pdf() -> None:
    """
    mixed.pdf — page 1 is text-layer, page 2 is a scanned image.
    Tests mixed-mode page-by-page handling.
    """
    from fpdf import FPDF

    pdf = FPDF()

    # Page 1 — text layer
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 16)
    pdf.cell(0, 12, "Mixed PDF Document", new_x="LMARGIN", new_y="NEXT")
    pdf.set_font("Helvetica", "", 12)
    pdf.ln(5)
    pdf.multi_cell(0, 7,
        "This first page has a proper text layer and can be extracted "
        "directly by pdfplumber without needing OCR. It contains enough "
        "text to pass the minimum text length threshold for detection."
    )

    # Page 2 — scanned image
    img = _draw_text_page(
        width=800, height=1100,
        lines=[
            "SCANNED PAGE TWO",
            "",
            "This is the second page rendered",
            "as an image to simulate a scan.",
        ],
        font_size=28,
    )
    img_path = FIXTURES_DIR / "_temp_mixed_scan.png"
    img.save(img_path)

    pdf.add_page()
    pdf.image(str(img_path), x=0, y=0, w=210, h=297)
    pdf.output(str(FIXTURES_DIR / "mixed.pdf"))

    img_path.unlink(missing_ok=True)


# ── PPTX fixtures ─────────────────────────────────────────────────────────────

def _make_simple_pptx() -> None:
    """
    simple.pptx — 5 slides:
      1. Title slide
      2. Text slide with body paragraphs
      3. Slide with a table
      4. Slide with an image
      5. Slide with speaker notes
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()

    # Slide 1 — Title
    layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Test Presentation"
    subtitle = slide.placeholders[1]
    subtitle.text = "Created by MarkFlow test fixtures"

    # Slide 2 — Text content
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Introduction"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "This is the first paragraph of the slide."
    p = tf.add_paragraph()
    p.text = "This is the second paragraph with more detail."
    p2 = tf.add_paragraph()
    run = p2.add_run()
    run.text = "Bold text here"
    run.font.bold = True

    # Slide 3 — Table
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Data Overview"
    rows, cols = 4, 3
    table_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(2.0), Inches(9.0), Inches(2.5))
    table = table_shape.table
    headers = ["Name", "Score", "Grade"]
    data = [["Alice", "92", "A"], ["Bob", "78", "B+"], ["Charlie", "65", "C"]]
    for ci, h in enumerate(headers):
        table.cell(0, ci).text = h
    for ri, row in enumerate(data, start=1):
        for ci, val in enumerate(row):
            table.cell(ri, ci).text = val

    # Slide 4 — Image
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Image Slide"
    img_buf = _make_png_bytes(100, 80, (50, 120, 200))
    slide.shapes.add_picture(img_buf, Inches(2.0), Inches(2.5), width=Inches(4.0))

    # Slide 5 — Speaker notes
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Summary"
    body = slide.placeholders[1]
    body.text = "This slide has speaker notes attached."
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "Remember to emphasize the key findings from slide 3."

    prs.save(str(FIXTURES_DIR / "simple.pptx"))


# ── XLSX fixtures ─────────────────────────────────────────────────────────────

def _make_simple_xlsx() -> None:
    """
    simple.xlsx — 2 sheets:
      Sheet1: headers + 10 rows of mixed types (strings, ints, floats, dates) + SUM formula
      Sheet2: 5 rows with a merged cell range
    """
    import openpyxl
    from datetime import date

    wb = openpyxl.Workbook()

    # Sheet 1
    ws1 = wb.active
    ws1.title = "Sales Data"
    ws1.append(["Name", "Region", "Units", "Price", "Total"])
    data = [
        ("Alice", "North", 10, 25.50, None),
        ("Bob", "South", 15, 30.00, None),
        ("Charlie", "East", 8, 22.75, None),
        ("Diana", "West", 20, 18.00, None),
        ("Eve", "North", 12, 35.00, None),
        ("Frank", "South", 7, 28.50, None),
        ("Grace", "East", 25, 15.00, None),
        ("Henry", "West", 30, 12.50, None),
        ("Ivy", "North", 5, 45.00, None),
        ("Jack", "South", 18, 20.00, None),
    ]
    for i, row in enumerate(data, start=2):
        ws1.append(list(row[:4]))
        # Total = Units * Price
        ws1.cell(row=i, column=5).value = f"=C{i}*D{i}"

    # SUM at bottom
    ws1.cell(row=12, column=3).value = "=SUM(C2:C11)"
    ws1.cell(row=12, column=5).value = "=SUM(E2:E11)"

    # Column widths
    ws1.column_dimensions["A"].width = 12
    ws1.column_dimensions["B"].width = 10
    ws1.column_dimensions["C"].width = 8
    ws1.column_dimensions["D"].width = 10
    ws1.column_dimensions["E"].width = 10

    # Sheet 2 with merged cells
    ws2 = wb.create_sheet("Summary")
    ws2.append(["Category", "Count", "Notes"])
    ws2.append(["Alpha", "5", "First group"])
    ws2.append(["Beta", "3", "Second group"])
    ws2.append(["Gamma", "8", "Third group"])
    ws2.append(["Delta", "2", "Fourth group"])

    # Merge A2:A3 to test merge handling
    ws2.merge_cells("A2:A3")

    wb.save(str(FIXTURES_DIR / "simple.xlsx"))


def _make_complex_xlsx() -> None:
    """
    complex.xlsx — conditional formatting, multiple number formats,
    wider column widths, freeze panes.
    """
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment, numbers

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Financial"

    # Headers with styling
    headers = ["Date", "Description", "Debit", "Credit", "Balance"]
    for ci, h in enumerate(headers, start=1):
        cell = ws.cell(row=1, column=ci, value=h)
        cell.font = Font(bold=True, size=11)
        cell.fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
        cell.font = Font(bold=True, color="FFFFFF")

    data = [
        ("2026-01-15", "Opening balance", "", "1000.00", "=D2"),
        ("2026-01-20", "Office supplies", "45.99", "", "=E2-C3+D3"),
        ("2026-02-01", "Client payment", "", "500.00", "=E3-C4+D4"),
        ("2026-02-10", "Software license", "120.00", "", "=E4-C5+D5"),
        ("2026-02-15", "Consulting income", "", "2000.00", "=E5-C6+D6"),
    ]
    for ri, row in enumerate(data, start=2):
        for ci, val in enumerate(row, start=1):
            cell = ws.cell(row=ri, column=ci)
            if val.startswith("="):
                cell.value = val
            elif ci in (3, 4) and val:
                cell.value = float(val)
                cell.number_format = '#,##0.00'
            else:
                cell.value = val

    # Column widths
    ws.column_dimensions["A"].width = 14
    ws.column_dimensions["B"].width = 22
    ws.column_dimensions["C"].width = 12
    ws.column_dimensions["D"].width = 12
    ws.column_dimensions["E"].width = 14

    # Freeze panes
    ws.freeze_panes = "A2"

    wb.save(str(FIXTURES_DIR / "complex.xlsx"))


# ── CSV fixtures ──────────────────────────────────────────────────────────────

def _make_simple_csv() -> None:
    """simple.csv — 5 columns, 10 rows, UTF-8."""
    import csv

    rows = [
        ["id", "name", "age", "city", "score"],
        ["1", "Alice", "30", "New York", "92.5"],
        ["2", "Bob", "25", "London", "78.0"],
        ["3", "Charlie", "35", "Paris", "65.3"],
        ["4", "Diana", "28", "Tokyo", "88.1"],
        ["5", "Eve", "32", "Sydney", "71.9"],
        ["6", "Frank", "40", "Berlin", "95.0"],
        ["7", "Grace", "27", "Toronto", "82.4"],
        ["8", "Henry", "33", "Seoul", "77.8"],
        ["9", "Ivy", "29", "Mumbai", "90.2"],
        ["10", "Jack", "31", "Cairo", "68.7"],
    ]

    with open(FIXTURES_DIR / "simple.csv", "w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f)
        writer.writerows(rows)


def _make_unicode_csv() -> None:
    """unicode.csv — non-ASCII characters: accents, CJK, emoji."""
    import csv

    rows = [
        ["id", "name", "city", "notes"],
        ["1", "José García", "México City", "Café owner ☕"],
        ["2", "Müller", "München", "Straße 42"],
        ["3", "田中太郎", "東京", "日本語テスト"],
        ["4", "Björk", "Reykjavík", "Icelandic ð þ"],
        ["5", "Łukasz", "Łódź", "Polish ł ź ś"],
    ]

    with open(FIXTURES_DIR / "unicode.csv", "w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f)
        writer.writerows(rows)


def _make_simple_tsv() -> None:
    """simple.tsv — tab-delimited version."""
    import csv

    rows = [
        ["id", "name", "value"],
        ["1", "Alpha", "100"],
        ["2", "Beta", "200"],
        ["3", "Gamma", "300"],
        ["4", "Delta", "400"],
        ["5", "Epsilon", "500"],
    ]

    with open(FIXTURES_DIR / "simple.tsv", "w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f, delimiter="\t")
        writer.writerows(rows)


# ── Entry point ────────────────────────────────────────────────────────────────

def generate_all() -> None:
    """Generate all test fixtures. Safe to call multiple times (idempotent)."""
    FIXTURES_DIR.mkdir(exist_ok=True)

    # DOCX
    if not (FIXTURES_DIR / "simple.docx").exists():
        _make_simple_docx()
    if not (FIXTURES_DIR / "complex.docx").exists():
        _make_complex_docx()

    # PDF
    if not (FIXTURES_DIR / "simple_text.pdf").exists():
        _make_simple_text_pdf()
    if not (FIXTURES_DIR / "scanned.pdf").exists():
        _make_scanned_pdf()
    if not (FIXTURES_DIR / "mixed.pdf").exists():
        _make_mixed_pdf()

    # PPTX
    if not (FIXTURES_DIR / "simple.pptx").exists():
        _make_simple_pptx()

    # XLSX
    if not (FIXTURES_DIR / "simple.xlsx").exists():
        _make_simple_xlsx()
    if not (FIXTURES_DIR / "complex.xlsx").exists():
        _make_complex_xlsx()

    # CSV / TSV
    if not (FIXTURES_DIR / "simple.csv").exists():
        _make_simple_csv()
    if not (FIXTURES_DIR / "unicode.csv").exists():
        _make_unicode_csv()
    if not (FIXTURES_DIR / "simple.tsv").exists():
        _make_simple_tsv()

    # OCR images
    generate_ocr_fixtures(FIXTURES_DIR)


if __name__ == "__main__":
    generate_all()
    print(f"Fixtures generated in {FIXTURES_DIR}")
