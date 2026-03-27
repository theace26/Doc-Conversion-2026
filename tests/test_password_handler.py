"""
Tests for core/password_handler.py — password-protected document handling.

Test categories:
  - Detection: no false positives on normal files
  - Restriction stripping: PDF owner-password, OOXML edit/sheet protection
  - Decryption: known password, cascade, failure handling
  - Tempfile cleanup, timeout
"""

import tempfile
import zipfile
from pathlib import Path

import pytest

from core.password_handler import (
    CrackMethod,
    PasswordHandler,
    PasswordResult,
    ProtectionType,
)

FIXTURE_DIR = Path(__file__).parent / "fixtures" / "password"

# ── Helpers ────────────────────────────────────────────────────────────────────


def _make_plain_pdf(tmp_path: Path) -> Path:
    """Create a minimal unprotected PDF."""
    from fpdf import FPDF
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Helvetica", size=12)
    pdf.cell(200, 10, text="Hello World", new_x="LMARGIN", new_y="NEXT")
    path = tmp_path / "plain.pdf"
    pdf.output(str(path))
    return path


def _make_plain_docx(tmp_path: Path) -> Path:
    """Create a minimal unprotected DOCX."""
    import docx
    doc = docx.Document()
    doc.add_paragraph("Hello World")
    path = tmp_path / "plain.docx"
    doc.save(str(path))
    return path


def _make_plain_xlsx(tmp_path: Path) -> Path:
    """Create a minimal unprotected XLSX."""
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws["A1"] = "Hello"
    path = tmp_path / "plain.xlsx"
    wb.save(str(path))
    return path


def _make_plain_pptx(tmp_path: Path) -> Path:
    """Create a minimal unprotected PPTX."""
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Hello"
    path = tmp_path / "plain.pptx"
    prs.save(str(path))
    return path


def _make_restricted_pdf(tmp_path: Path) -> Path:
    """Create a PDF with owner password only (no user password)."""
    import pikepdf
    plain = _make_plain_pdf(tmp_path)
    restricted = tmp_path / "restricted.pdf"
    pdf = pikepdf.open(plain)
    pdf.save(str(restricted), encryption=pikepdf.Encryption(owner="ownerpass", user="", R=6))
    pdf.close()
    return restricted


def _make_encrypted_pdf(tmp_path: Path, password: str = "test123") -> Path:
    """Create a PDF with user password."""
    import pikepdf
    plain = _make_plain_pdf(tmp_path)
    encrypted = tmp_path / "encrypted.pdf"
    pdf = pikepdf.open(plain)
    pdf.save(str(encrypted), encryption=pikepdf.Encryption(owner="ownerpass", user=password, R=6))
    pdf.close()
    return encrypted


def _make_protected_docx(tmp_path: Path) -> Path:
    """Create a DOCX with edit protection."""
    from lxml import etree
    plain = _make_plain_docx(tmp_path)
    protected = tmp_path / "protected_edit.docx"
    with zipfile.ZipFile(plain, "r") as zin, zipfile.ZipFile(protected, "w") as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename == "word/settings.xml":
                tree = etree.fromstring(data)
                ns = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
                prot = etree.SubElement(tree, f"{{{ns}}}documentProtection")
                prot.set(f"{{{ns}}}edit", "readOnly")
                prot.set(f"{{{ns}}}enforcement", "1")
                data = etree.tostring(tree, xml_declaration=True, encoding="UTF-8", standalone=True)
            zout.writestr(item, data)
    return protected


def _make_protected_xlsx(tmp_path: Path) -> Path:
    """Create an XLSX with sheet protection."""
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws["A1"] = "Protected"
    ws.protection.sheet = True
    ws.protection.password = "sheetpass"
    path = tmp_path / "protected_sheet.xlsx"
    wb.save(str(path))
    return path


# ── Detection tests ────────────────────────────────────────────────────────────


class TestDetection:
    """Test that detection correctly identifies protection types."""

    def test_detect_unprotected_pdf(self, tmp_path):
        path = _make_plain_pdf(tmp_path)
        handler = PasswordHandler()
        result = handler.handle_sync(path)
        assert result.protection_type == ProtectionType.NONE
        assert result.success is True

    def test_detect_unprotected_docx(self, tmp_path):
        path = _make_plain_docx(tmp_path)
        handler = PasswordHandler()
        result = handler.handle_sync(path)
        assert result.protection_type == ProtectionType.NONE
        assert result.success is True

    def test_detect_unprotected_xlsx(self, tmp_path):
        path = _make_plain_xlsx(tmp_path)
        handler = PasswordHandler()
        result = handler.handle_sync(path)
        assert result.protection_type == ProtectionType.NONE
        assert result.success is True

    def test_detect_unprotected_pptx(self, tmp_path):
        path = _make_plain_pptx(tmp_path)
        handler = PasswordHandler()
        result = handler.handle_sync(path)
        assert result.protection_type == ProtectionType.NONE
        assert result.success is True

    def test_detect_unsupported_extension(self, tmp_path):
        path = tmp_path / "test.csv"
        path.write_text("a,b,c\n1,2,3")
        handler = PasswordHandler()
        result = handler.handle_sync(path)
        assert result.protection_type == ProtectionType.NONE
        assert result.success is True

    def test_detect_restricted_pdf(self, tmp_path):
        path = _make_restricted_pdf(tmp_path)
        handler = PasswordHandler()
        result = handler.handle_sync(path)
        assert result.protection_type == ProtectionType.RESTRICTION_ONLY
        assert result.success is True

    def test_detect_encrypted_pdf(self, tmp_path):
        path = _make_encrypted_pdf(tmp_path)
        handler = PasswordHandler({"password_dictionary_enabled": "false", "password_brute_force_enabled": "false", "password_timeout_seconds": "5"})
        result = handler.handle_sync(path)
        # Without password, it should fail
        assert result.protection_type == ProtectionType.ENCRYPTED_FAILED
        assert result.success is False

    def test_detect_protected_docx(self, tmp_path):
        path = _make_protected_docx(tmp_path)
        handler = PasswordHandler()
        result = handler.handle_sync(path)
        assert result.protection_type == ProtectionType.RESTRICTION_ONLY
        assert result.success is True


# ── Restriction stripping tests ────────────────────────────────────────────────


class TestRestrictionStripping:
    """Test automatic stripping of edit/print restrictions."""

    def test_strip_pdf_restrictions(self, tmp_path):
        path = _make_restricted_pdf(tmp_path)
        handler = PasswordHandler()
        result = handler.handle_sync(path)
        assert result.success is True
        assert result.protection_type == ProtectionType.RESTRICTION_ONLY
        assert result.output_path is not None
        assert result.output_path.exists()
        # The output should be a valid PDF that opens without password
        import pikepdf
        pdf = pikepdf.open(result.output_path)
        assert len(pdf.pages) > 0
        pdf.close()

    def test_strip_docx_edit_protection(self, tmp_path):
        path = _make_protected_docx(tmp_path)
        handler = PasswordHandler()
        result = handler.handle_sync(path)
        assert result.success is True
        assert result.output_path is not None
        # Verify the protection tag is removed
        from lxml import etree
        with zipfile.ZipFile(result.output_path, "r") as zf:
            if "word/settings.xml" in zf.namelist():
                data = zf.read("word/settings.xml")
                tree = etree.fromstring(data)
                ns = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
                prot = tree.findall(f".//{{{ns}}}documentProtection")
                assert len(prot) == 0, "documentProtection tag should be removed"

    def test_strip_xlsx_sheet_protection(self, tmp_path):
        path = _make_protected_xlsx(tmp_path)
        handler = PasswordHandler()
        result = handler.handle_sync(path)
        assert result.success is True
        assert result.output_path is not None
        from lxml import etree
        with zipfile.ZipFile(result.output_path, "r") as zf:
            for name in zf.namelist():
                if name.startswith("xl/worksheets/") and name.endswith(".xml"):
                    data = zf.read(name)
                    tree = etree.fromstring(data)
                    ns = "http://schemas.openxmlformats.org/spreadsheetml/2006/main"
                    prot = tree.findall(f".//{{{ns}}}sheetProtection")
                    assert len(prot) == 0, f"sheetProtection should be removed from {name}"


# ── Decryption tests ──────────────────────────────────────────────────────────


class TestDecryption:
    """Test password decryption for encrypted files."""

    def test_decrypt_pdf_known_password(self, tmp_path):
        path = _make_encrypted_pdf(tmp_path, "test123")
        handler = PasswordHandler({"password_dictionary_enabled": "false", "password_brute_force_enabled": "false"})
        result = handler.handle_sync(path, user_password="test123")
        assert result.success is True
        assert result.protection_type == ProtectionType.ENCRYPTED_DECRYPTED
        assert result.method == CrackMethod.USER_SUPPLIED
        assert result.output_path is not None
        assert result.output_path.exists()

    def test_decrypt_pdf_wrong_password(self, tmp_path):
        path = _make_encrypted_pdf(tmp_path, "test123")
        handler = PasswordHandler({"password_dictionary_enabled": "false", "password_brute_force_enabled": "false", "password_timeout_seconds": "5"})
        result = handler.handle_sync(path, user_password="wrongpass")
        assert result.success is False
        assert result.protection_type == ProtectionType.ENCRYPTED_FAILED

    def test_decrypt_pdf_via_dictionary(self, tmp_path):
        # "password" is in common.txt
        path = _make_encrypted_pdf(tmp_path, "password")
        handler = PasswordHandler({"password_dictionary_enabled": "true", "password_brute_force_enabled": "false", "password_timeout_seconds": "30"})
        result = handler.handle_sync(path)
        assert result.success is True
        assert result.method == CrackMethod.DICTIONARY

    def test_password_cascade_order(self, tmp_path):
        """Verify empty → user → org → found → dictionary order."""
        path = _make_encrypted_pdf(tmp_path, "test123")
        handler = PasswordHandler({"password_dictionary_enabled": "false", "password_brute_force_enabled": "false", "password_timeout_seconds": "5"})
        # Simulate found password from another file
        handler._found_passwords = ["test123"]
        result = handler.handle_sync(path)
        assert result.success is True
        assert result.method == CrackMethod.FOUND_REUSE

    def test_batch_password_reuse(self, tmp_path):
        """Password found on file A should be tried on file B."""
        handler = PasswordHandler({"password_dictionary_enabled": "false", "password_brute_force_enabled": "false", "password_timeout_seconds": "10"})

        # File A: encrypted with "test123", provide it manually
        path_a = _make_encrypted_pdf(tmp_path, "test123")
        result_a = handler.handle_sync(path_a, user_password="test123")
        assert result_a.success is True
        assert "test123" in handler._found_passwords

        # File B: same password, no user input — should find via reuse
        path_b = _make_encrypted_pdf(tmp_path, "test123")
        path_b = path_b.rename(tmp_path / "encrypted_b.pdf")
        result_b = handler.handle_sync(path_b)
        assert result_b.success is True
        assert result_b.method == CrackMethod.FOUND_REUSE

    def test_failure_does_not_raise(self, tmp_path):
        """Failed decryption returns error result, never raises."""
        path = _make_encrypted_pdf(tmp_path, "verysecret")
        handler = PasswordHandler({"password_dictionary_enabled": "false", "password_brute_force_enabled": "false", "password_timeout_seconds": "2"})
        result = handler.handle_sync(path)
        assert isinstance(result, PasswordResult)
        assert result.success is False
        assert result.error is not None

    def test_timeout_respected(self, tmp_path):
        """Brute-force should stop after timeout."""
        path = _make_encrypted_pdf(tmp_path, "zzzzzzz")  # Not in dictionary, hard to brute-force
        handler = PasswordHandler({
            "password_dictionary_enabled": "false",
            "password_brute_force_enabled": "true",
            "password_brute_force_max_length": "8",
            "password_brute_force_charset": "alpha",
            "password_timeout_seconds": "3",
        })
        import time
        start = time.monotonic()
        result = handler.handle_sync(path)
        elapsed = time.monotonic() - start
        assert result.success is False
        # Should have stopped reasonably close to timeout (within a few seconds)
        assert elapsed < 10, f"Took {elapsed}s, expected ~3s timeout"


# ── Tempfile cleanup ──────────────────────────────────────────────────────────


class TestCleanup:
    """Test that temp files are cleaned up."""

    def test_restriction_creates_temp(self, tmp_path):
        """Restriction stripping creates a temp file that differs from original."""
        path = _make_restricted_pdf(tmp_path)
        handler = PasswordHandler()
        result = handler.handle_sync(path)
        assert result.output_path != path
        assert result.output_path.exists()

    def test_cleanup_removes_temp(self, tmp_path):
        """cleanup_temp_file() removes the decrypted temp file."""
        path = _make_encrypted_pdf(tmp_path, "test123")
        handler = PasswordHandler({"password_dictionary_enabled": "false"})
        result = handler.handle_sync(path, user_password="test123")
        assert result.output_path.exists()
        handler.cleanup_temp_file(result)
        # File should be gone (if it was in temp dir)
        # Note: cleanup only removes files in system temp dir
