"""
Generate password-protected test fixtures for all supported formats.

Run inside Docker: python tests/generate_password_fixtures.py
Requires: pikepdf, msoffcrypto-tool, python-docx, openpyxl, python-pptx
"""

import os
import sys
from pathlib import Path

FIXTURE_DIR = Path(__file__).parent / "fixtures" / "password"
FIXTURE_DIR.mkdir(parents=True, exist_ok=True)

TEST_PASSWORD = "test123"


def create_plain_pdf(path: Path) -> None:
    """Create a simple PDF file."""
    from fpdf import FPDF
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Helvetica", size=12)
    pdf.cell(200, 10, text="This is a test PDF document.", new_x="LMARGIN", new_y="NEXT")
    pdf.cell(200, 10, text="It has two lines of text.", new_x="LMARGIN", new_y="NEXT")
    pdf.output(str(path))


def create_restricted_pdf(path: Path) -> None:
    """Create a PDF with owner password (restrictions only, no user password)."""
    import pikepdf
    plain = FIXTURE_DIR / "_temp_plain.pdf"
    create_plain_pdf(plain)
    pdf = pikepdf.open(plain)
    pdf.save(str(path), encryption=pikepdf.Encryption(
        owner="ownerpass",
        user="",  # empty user password = can open freely
        R=6,
    ))
    pdf.close()
    plain.unlink()


def create_encrypted_pdf(path: Path) -> None:
    """Create a PDF encrypted with user password."""
    import pikepdf
    plain = FIXTURE_DIR / "_temp_plain2.pdf"
    create_plain_pdf(plain)
    pdf = pikepdf.open(plain)
    pdf.save(str(path), encryption=pikepdf.Encryption(
        owner="ownerpass",
        user=TEST_PASSWORD,
        R=6,
    ))
    pdf.close()
    plain.unlink()


def create_plain_docx(path: Path) -> None:
    """Create a simple DOCX file."""
    import docx
    doc = docx.Document()
    doc.add_heading("Test Document", 0)
    doc.add_paragraph("This is a test paragraph.")
    doc.save(str(path))


def create_protected_edit_docx(path: Path) -> None:
    """Create a DOCX with edit protection (no encryption)."""
    import docx
    from lxml import etree

    doc = docx.Document()
    doc.add_heading("Protected Document", 0)
    doc.add_paragraph("This document has edit protection enabled.")
    # Save first
    temp = FIXTURE_DIR / "_temp_docx_edit.docx"
    doc.save(str(temp))

    # Add documentProtection tag to settings.xml
    import zipfile
    with zipfile.ZipFile(temp, "r") as zin, zipfile.ZipFile(path, "w") as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename == "word/settings.xml":
                tree = etree.fromstring(data)
                nsmap = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
                prot = etree.SubElement(tree, f"{{{nsmap['w']}}}documentProtection")
                prot.set(f"{{{nsmap['w']}}}edit", "readOnly")
                prot.set(f"{{{nsmap['w']}}}enforcement", "1")
                data = etree.tostring(tree, xml_declaration=True, encoding="UTF-8", standalone=True)
            zout.writestr(item, data)
    temp.unlink()


def create_encrypted_docx(path: Path) -> None:
    """Create an encrypted DOCX."""
    import msoffcrypto
    plain = FIXTURE_DIR / "_temp_docx_enc.docx"
    create_plain_docx(plain)
    with open(plain, "rb") as f:
        file = msoffcrypto.OfficeFile(f)
        file.load_key(password=TEST_PASSWORD)
        with open(path, "wb") as out:
            file.encrypt(out)  # msoffcrypto.encrypt for creating encrypted files
    plain.unlink()


def create_protected_sheet_xlsx(path: Path) -> None:
    """Create an XLSX with sheet protection."""
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws["A1"] = "Protected Cell"
    ws.protection.sheet = True
    ws.protection.password = "sheetpass"
    wb.save(str(path))


def create_encrypted_xlsx(path: Path) -> None:
    """Create an encrypted XLSX."""
    import msoffcrypto
    import openpyxl
    plain = FIXTURE_DIR / "_temp_xlsx_enc.xlsx"
    wb = openpyxl.Workbook()
    ws = wb.active
    ws["A1"] = "Encrypted Data"
    wb.save(str(plain))
    with open(plain, "rb") as f:
        file = msoffcrypto.OfficeFile(f)
        file.load_key(password=TEST_PASSWORD)
        with open(path, "wb") as out:
            file.encrypt(out)
    plain.unlink()


def create_protected_pptx(path: Path) -> None:
    """Create a PPTX with modify protection."""
    from pptx import Presentation
    from lxml import etree
    import zipfile

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Protected Presentation"
    temp = FIXTURE_DIR / "_temp_pptx_prot.pptx"
    prs.save(str(temp))

    # Add modifyVerifier tag
    with zipfile.ZipFile(temp, "r") as zin, zipfile.ZipFile(path, "w") as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename == "ppt/presentation.xml":
                tree = etree.fromstring(data)
                ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
                mv = etree.SubElement(tree, f"{{{ns}}}modifyVerifier")
                mv.set("cryptProviderType", "rsaAES")
                mv.set("hashData", "dummyhash")
                data = etree.tostring(tree, xml_declaration=True, encoding="UTF-8", standalone=True)
            zout.writestr(item, data)
    temp.unlink()


def create_encrypted_pptx(path: Path) -> None:
    """Create an encrypted PPTX."""
    import msoffcrypto
    from pptx import Presentation

    plain = FIXTURE_DIR / "_temp_pptx_enc.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Encrypted Presentation"
    prs.save(str(plain))

    with open(plain, "rb") as f:
        file = msoffcrypto.OfficeFile(f)
        file.load_key(password=TEST_PASSWORD)
        with open(path, "wb") as out:
            file.encrypt(out)
    plain.unlink()


def main():
    print(f"Generating fixtures in {FIXTURE_DIR}")

    # PDF
    create_restricted_pdf(FIXTURE_DIR / "restricted.pdf")
    print("  restricted.pdf")
    create_encrypted_pdf(FIXTURE_DIR / "encrypted.pdf")
    print("  encrypted.pdf")

    # DOCX
    create_protected_edit_docx(FIXTURE_DIR / "protected_edit.docx")
    print("  protected_edit.docx")
    try:
        create_encrypted_docx(FIXTURE_DIR / "encrypted.docx")
        print("  encrypted.docx")
    except Exception as e:
        print(f"  encrypted.docx SKIPPED: {e}")

    # XLSX
    create_protected_sheet_xlsx(FIXTURE_DIR / "protected_sheet.xlsx")
    print("  protected_sheet.xlsx")
    try:
        create_encrypted_xlsx(FIXTURE_DIR / "encrypted.xlsx")
        print("  encrypted.xlsx")
    except Exception as e:
        print(f"  encrypted.xlsx SKIPPED: {e}")

    # PPTX
    create_protected_pptx(FIXTURE_DIR / "protected.pptx")
    print("  protected.pptx")
    try:
        create_encrypted_pptx(FIXTURE_DIR / "encrypted.pptx")
        print("  encrypted.pptx")
    except Exception as e:
        print(f"  encrypted.pptx SKIPPED: {e}")

    print("Done!")


if __name__ == "__main__":
    main()
